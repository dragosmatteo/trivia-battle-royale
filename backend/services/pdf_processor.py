"""
PDF, Word (.docx), and PowerPoint (.pptx) text extraction with Romanian diacritics repair.
Handles LaTeX-generated PDFs where diacritics appear as combining characters.
"""

import os
import re

# Try pdfplumber first (better extraction), fall back to PyPDF2
try:
    import pdfplumber
    _USE_PDFPLUMBER = True
except ImportError:
    _USE_PDFPLUMBER = False

from PyPDF2 import PdfReader


def extract_text_from_pdf(file_path: str) -> str:
    """Extract and clean text from a PDF file with Romanian diacritics repair."""
    raw_text = ""

    if _USE_PDFPLUMBER:
        try:
            with pdfplumber.open(file_path) as pdf:
                for page in pdf.pages:
                    page_text = page.extract_text() or ""
                    if page_text:
                        raw_text += page_text + "\n"
        except Exception:
            raw_text = ""

    # Fallback to PyPDF2 if pdfplumber failed or isn't available
    if not raw_text.strip():
        try:
            reader = PdfReader(file_path)
            for page in reader.pages:
                page_text = page.extract_text() or ""
                if page_text:
                    raw_text += page_text + "\n"
        except Exception:
            pass

    if not raw_text.strip():
        return ""

    # Step 1: Fix Romanian diacritics from LaTeX PDFs
    text = _fix_romanian_diacritics(raw_text)

    # Step 2: Fix concatenated words (LaTeX table-of-contents artifacts)
    text = _fix_concatenated_words(text)

    # Step 3: General cleanup
    text = _clean_text(text)

    return text


def _fix_romanian_diacritics(text: str) -> str:
    """
    Repair Romanian diacritics broken by LaTeX PDF generation.

    LaTeX PDFs produce combining marks on separate lines:
      a˘ or ˘a  → ă       (a-breve)
      standalone , on next line → artifact from ș/ț comma-below
    """
    # --- Step 0: Remove standalone comma lines (LaTeX comma-below artifacts) ---
    # These are orphan marks that can't be reliably mapped back to letters.
    # Pattern: lines containing only commas/spaces like ",", ", ," etc.
    text = re.sub(r'^\s*[,\s]+\s*$', '', text, flags=re.MULTILINE)

    # --- Step 1: Fix breve accent (˘) combinations ---
    # Pattern: letter followed by breve (e.g., "a˘" → "ă")
    text = re.sub(r'a˘', 'ă', text)
    text = re.sub(r'A˘', 'Ă', text)
    # Pattern: breve followed by letter (e.g., "˘a" → "ă")
    text = re.sub(r'˘a', 'ă', text)
    text = re.sub(r'˘A', 'Ă', text)
    # Handle breve with newline between: "a\n˘" or "˘\na"
    text = re.sub(r'a\s*\n\s*˘', 'ă', text)
    text = re.sub(r'˘\s*\n\s*a', 'ă', text)

    # Breve with i (for î alternate representation in some LaTeX)
    text = re.sub(r'ˆı', 'î', text)
    text = re.sub(r'ˆI', 'Î', text)

    # --- Step 2: Fix cedilla variants (older Unicode for ș/ț) ---
    text = text.replace('ţ', 'ț')
    text = text.replace('Ţ', 'Ț')
    text = text.replace('ş', 'ș')
    text = text.replace('Ş', 'Ș')

    # --- Step 3: Fix inline comma-below patterns ---
    # "s," or "t," when followed by a letter → ș/ț
    text = re.sub(r's,(?=[a-zA-ZăîâĂÎÂ])', 'ș', text)
    text = re.sub(r'S,(?=[a-zA-ZăîâĂÎÂ])', 'Ș', text)
    text = re.sub(r't,(?=[a-zA-ZăîâĂÎÂ])', 'ț', text)
    text = re.sub(r'T,(?=[a-zA-ZăîâĂÎÂ])', 'Ț', text)

    # --- Step 4: Remove orphaned breve marks ---
    text = re.sub(r'^\s*˘\s*$', '', text, flags=re.MULTILINE)
    text = re.sub(r'(?<!\w)˘(?!\w)', '', text)

    # --- Step 5: Common Romanian word corrections ---
    # Fix patterns that result from imperfect diacritics recovery
    common_fixes = {
        'educationala': 'educațională',
        'aplicatie': 'aplicație',
        'aplicatii': 'aplicații',
        'traditionale': 'tradiționale',
        'informatie': 'informație',
        'informatii': 'informații',
        'stiintific': 'științific',
        'competitie': 'competiție',
        'evaluare': 'evaluare',
        'functionala': 'funcțională',
        'functionale': 'funcționale',
        'utilizatorilor': 'utilizatorilor',
        'implementare': 'implementare',
        'interactiune': 'interacțiune',
        'comunicatie': 'comunicație',
        'sectiune': 'secțiune',
        'directie': 'direcție',
        'autentificare': 'autentificare',
        'sesiune': 'sesiune',
    }
    for wrong, right in common_fixes.items():
        text = re.sub(re.escape(wrong), right, text, flags=re.IGNORECASE)

    return text


def _fix_concatenated_words(text: str) -> str:
    """
    Fix words concatenated by LaTeX (e.g., 'ContextulSchimbării' → 'Contextul Schimbării').
    Only split on CamelCase patterns where a lowercase letter is followed by uppercase.
    """
    # Insert space between lowercase and uppercase (CamelCase split)
    text = re.sub(r'([a-zăîâșț])([A-ZĂÎÂȘȚ])', r'\1 \2', text)
    # Fix double/triple spaces
    text = re.sub(r'  +', ' ', text)
    return text


def _clean_text(text: str) -> str:
    """Remove headers, footers, page numbers, and irrelevant artifacts."""
    lines = text.split("\n")
    cleaned_lines = []

    for line in lines:
        stripped = line.strip()
        # Skip empty lines
        if not stripped:
            continue
        # Skip page numbers (standalone numbers)
        if re.match(r'^\d{1,3}$', stripped):
            continue
        # Skip very short lines that are likely headers/footers/artifacts
        if len(stripped) < 3:
            continue
        # Skip lines that are just dots (table of contents leaders)
        if re.match(r'^[.\s]+$', stripped):
            continue
        cleaned_lines.append(stripped)

    return "\n".join(cleaned_lines)


def extract_text_from_docx(file_path: str) -> str:
    """Extract text from a Word (.docx) file using python-docx."""
    try:
        from docx import Document
        doc = Document(file_path)
        paragraphs = [para.text for para in doc.paragraphs if para.text.strip()]
        return "\n".join(paragraphs)
    except Exception:
        return ""


def extract_text_from_pptx(file_path: str) -> str:
    """Extract text from a PowerPoint (.pptx) file using python-pptx."""
    try:
        from pptx import Presentation
        prs = Presentation(file_path)
        lines = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if text:
                            lines.append(text)
        return "\n".join(lines)
    except Exception:
        return ""


def extract_text_from_file(file_path: str) -> str:
    """Detect file extension and extract text using the appropriate extractor."""
    ext = os.path.splitext(file_path)[1].lower()
    if ext == ".pdf":
        return extract_text_from_pdf(file_path)
    elif ext == ".docx":
        return extract_text_from_docx(file_path)
    elif ext == ".pptx":
        return extract_text_from_pptx(file_path)
    else:
        return ""


def chunk_text(text: str, chunk_size: int = 2000, overlap: int = 200) -> list[str]:
    """Split text into overlapping chunks for RAG processing."""
    words = text.split()
    if not words:
        return []

    chunks = []
    start = 0
    words_per_chunk = max(chunk_size // 5, 50)  # approximate words per chunk
    overlap_words = max(overlap // 5, 10)

    while start < len(words):
        end = min(start + words_per_chunk, len(words))
        chunk = " ".join(words[start:end])
        chunks.append(chunk)

        if end >= len(words):
            break  # Last chunk reached

        start = end - overlap_words  # overlap for context continuity

    return chunks
